import sys
import json
from docx import Document
from pptx import Presentation
import re
from multiprocessing import Pool
import os
from openpyxl import load_workbook
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity

def extract_word_text(filepath):
    """Extract text from Word document"""
    try:
        doc = Document(filepath)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + " "
        return text.strip()
    except Exception as e:
        return None

def load_excel_fast(filepath):
    """
    Fast Excel loading using openpyxl directly
    Returns dict of sheet_name -> list of rows (as lists)
    """
    try:
        # read_only=True and data_only=True make it MUCH faster
        wb = load_workbook(filepath, read_only=True, data_only=True)
        
        sheets_data = {}
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            
            # Read all rows at once (generator for memory efficiency)
            rows = []
            for row in ws.iter_rows(values_only=True):
                # Skip completely empty rows
                if any(cell is not None for cell in row):
                    rows.append(row)
            
            sheets_data[sheet_name] = rows
        
        wb.close()
        return sheets_data
        
    except Exception as e:
        print(f"Error loading Excel {filepath}: {e}", file=sys.stderr)
        return None

def compare_excel_fast(data1, data2):
    """
    Compare Excel data loaded with openpyxl (fast!)
    data1, data2 = dict of sheet_name -> list of rows
    """
    try:
        if data1 is None or data2 is None:
            return 0.0
        
        # Get common sheet names
        common_sheets = set(data1.keys()).intersection(set(data2.keys()))
        
        if not common_sheets:
            return 0.0
        
        total_similarity = 0
        sheet_count = 0
        
        for sheet_name in common_sheets:
            sheet_similarity = compare_sheets_fast(data1[sheet_name], data2[sheet_name])
            total_similarity += sheet_similarity
            sheet_count += 1
        
        if sheet_count == 0:
            return 0.0
        
        return total_similarity / sheet_count
        
    except Exception as e:
        print(f"Error comparing Excel: {e}", file=sys.stderr)
        return 0.0

def compare_sheets_fast(rows1, rows2):
    """
    Compare two sheets (as lists of rows)
    Much faster than pandas approach!
    """
    if not rows1 or not rows2:
        return 0.0
    
    # Compare overlapping region
    min_rows = min(len(rows1), len(rows2))
    
    matching_cells = 0
    compared_cells = 0
    
    for i in range(min_rows):
        row1 = rows1[i]
        row2 = rows2[i]
        
        # Compare up to minimum column count
        min_cols = min(len(row1), len(row2))
        
        for j in range(min_cols):
            val1 = row1[j]
            val2 = row2[j]
            
            compared_cells += 1
            
            # Both None/empty
            if val1 is None and val2 is None:
                matching_cells += 1
            # Both have same value
            elif val1 == val2:
                matching_cells += 1
            # Try string comparison (for numbers vs strings)
            elif str(val1) == str(val2):
                matching_cells += 1
    
    return matching_cells / compared_cells if compared_cells > 0 else 0.0

def extract_powerpoint_text(filepath):
    """Extract text from PowerPoint"""
    try:
        prs = Presentation(filepath)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + " "
        return text.strip()
    except Exception as e:
        return None

def calculate_text_similarity(text1, text2):
    """Fast TF-IDF similarity"""
    if not text1 or not text2:
        return 0.0
    
    try:
        vectorizer = TfidfVectorizer()
        tfidf_matrix = vectorizer.fit_transform([text1, text2])
        similarity = cosine_similarity(tfidf_matrix[0:1], tfidf_matrix[1:2])[0][0]
        return float(similarity)
    except:
        return 0.0

def compare_files_batch(comparisons):
    """
    Compare multiple file pairs at once - PARALLEL VERSION
    """
    # Step 1: Collect all unique files that need to be read
    files_to_read = {}
    for comp in comparisons:
        file_type = comp['type']
        file1 = comp['file1']
        file2 = comp['file2']
        
        if file1 not in files_to_read:
            files_to_read[file1] = file_type
        if file2 not in files_to_read:
            files_to_read[file2] = file_type
    
    # Step 2: Read all files in PARALLEL (this is the slow part!)
    cpu_count = max(1, os.cpu_count() - 1)  # Leave 1 core free
    with Pool(processes=cpu_count) as pool:
        file_list = list(files_to_read.items())
        loaded_data = pool.starmap(load_file, file_list)
    
    # Create dictionary: filepath -> content
    file_cache = dict(zip(files_to_read.keys(), loaded_data))
    
    # Step 3: Compare using cached data (super fast!)
    results = {}
    
    for i, comp in enumerate(comparisons):
        file_type = comp['type']
        file1 = comp['file1']
        file2 = comp['file2']
        
        data1 = file_cache[file1]
        data2 = file_cache[file2]
        
        # Handle Excel files
        if file_type == 'excel':
            if data1 is None or data2 is None:
                similarity = 0.0
            else:
                similarity = compare_excel_fast(data1, data2)
            similar = similarity > 0.7
            results[str(i)] = {'similar': similar, 'score': similarity if similar else 0.0}
        
        # Handle Word files
        elif file_type == 'word':
            if data1 is None or data2 is None:
                results[str(i)] = {'similar': False, 'score': 0.0}
            else:
                similarity = calculate_text_similarity(data1, data2)
                similar = similarity > 0.6
                results[str(i)] = {'similar': similar, 'score': similarity if similar else 0.0}
                
        # Handle PowerPoint files
        elif file_type == 'powerpoint':
            if data1 is None or data2 is None:
                results[str(i)] = {'similar': False, 'score': 0.0}
            else:
                similarity = calculate_text_similarity(data1, data2)
                similar = similarity > 0.6
                results[str(i)] = {'similar': similar, 'score': similarity if similar else 0.0}
    
    return results


def load_file(filepath, file_type):
    """
    Load a single file - this runs in parallel!
    """
    try:
        if file_type == 'excel':
            return load_excel_fast(filepath)
        elif file_type == 'word':
            return extract_word_text(filepath)
        elif file_type == 'powerpoint':
            return extract_powerpoint_text(filepath)
    except Exception as e:
        print(f"Error loading {filepath}: {e}", file=sys.stderr)
        return None


if __name__ == "__main__":
    
    # Read JSON from stdin
    input_data = sys.stdin.read()
    
    try:
        comparisons = json.loads(input_data)
        results = compare_files_batch(comparisons)
        
        # Output results as JSON
        print(json.dumps(results))
        sys.exit(0)
    except Exception as e:
        print(json.dumps({'error': str(e)}), file=sys.stderr)
        sys.exit(1)