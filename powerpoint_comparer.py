import sys
import os
from pptx import Presentation
import re

def compare_pptx_files(file1, file2):
    """
    Compare two PowerPoint files for similarity based on text content.
    Returns True if files are similar (similarity > 0.6), False otherwise.
    """
    try:
        # DEBUG: Log file names
        with open("powerpoint_debug.log", "a", encoding="utf-8") as f:
            f.write(f"=== Comparing {os.path.basename(file1)} vs {os.path.basename(file2)} ===\n")
        
        # Read PowerPoint presentations using python-pptx library
        prs1 = Presentation(file1)
        prs2 = Presentation(file2)
        
        # Extract text content from both presentations
        text1 = extract_text(prs1)
        text2 = extract_text(prs2)
        
        # DEBUG: Log extracted texts
        with open("powerpoint_debug.log", "a", encoding="utf-8") as f:
            f.write(f"Text1 length: {len(text1)}, Text2 length: {len(text2)}\n")
            f.write(f"Text1: '{text1[:200]}'\n")
            f.write(f"Text2: '{text2[:200]}'\n")
        
        # Calculate similarity score between extracted texts
        similarity = calculate_text_similarity(text1, text2)
        
        # DEBUG: Log result
        with open("powerpoint_debug.log", "a", encoding="utf-8") as f:
            f.write(f"Similarity: {similarity}\n")
            f.write(f"Result: {'SIMILAR' if similarity > 0.6 else 'DIFFERENT'}\n\n")
        
        # Consider files similar if similarity score exceeds 60%
        return similarity > 0.6
        
    except Exception as e:
        # Log error for debugging
        with open("powerpoint_error.log", "a", encoding="utf-8") as f:
            f.write(f"Error: {str(e)}\n")
            f.write(f"File1: {file1}\n")
            f.write(f"File2: {file2}\n\n")
        return False

def extract_text(prs):
    """
    Extract all text content from a PowerPoint presentation.
    
    Args:
        prs: PowerPoint Presentation object
        
    Returns:
        String containing all text from slides and shapes
    """
    text = ""
    # Iterate through all slides in the presentation
    for slide in prs.slides:
        # Iterate through all shapes in each slide
        for shape in slide.shapes:
            # Check if shape has text property and extract text
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()

def calculate_text_similarity(text1, text2):
    """
    Calculate similarity between two text strings using character-based comparison.
    This preserves Umlaute and works better for German text.
    
    Args:
        text1: First text string
        text2: Second text string
        
    Returns:
        Similarity score between 0.0 (completely different) and 1.0 (identical)
    """
    # Return 0 if either text is empty
    if not text1 or not text2:
        return 0.0
    
    # For very short texts, use exact matching
    if len(text1) < 10 or len(text2) < 10:
        return 1.0 if text1 == text2 else 0.0
    
    # Character-based similarity (preserves Umlaute)
    common_chars = 0
    min_length = min(len(text1), len(text2))
    
    for i in range(min_length):
        if text1[i] == text2[i]:
            common_chars += 1
    
    similarity = common_chars / max(len(text1), len(text2))
    
    return similarity

if __name__ == "__main__":
    # Command line interface for standalone execution
    if len(sys.argv) != 3:
        print("Usage: python powerpoint_comparer.py <file1> <file2>")
        sys.exit(1)
    
    # Get file paths from command line arguments
    file1, file2 = sys.argv[1], sys.argv[2]
    
    # Compare files and exit with appropriate code
    # Exit code 0 = similar, 1 = not similar
    similar = compare_pptx_files(file1, file2)
    sys.exit(0 if similar else 1)